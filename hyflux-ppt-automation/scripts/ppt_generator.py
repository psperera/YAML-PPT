#!/usr/bin/env python3
"""
HyFlux PowerPoint Generator
Automates presentation creation from YAML content specs.
"""

import sys
import yaml
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import re

# Layout mapping (index to friendly name)
LAYOUT_MAP = {
    'title_white': 0,
    'title_reverse': 1,
    'divider': 6,
    'text_only': 8,
    'title_only': 12,
    'text_content': 21,
    'two_column': 23,
    'three_column': 24,
    'quote': 34,
    'end_slide': 35
}

class HyFluxPPTGenerator:
    def __init__(self, template_path, config_path=None):
        """Initialize generator with template."""
        self.template_path = Path(template_path)
        self.prs = Presentation(str(self.template_path))
        self.config = self._load_config(config_path)
        
        # Validate template
        if len(self.prs.slide_layouts) < 36:
            raise ValueError(f"Template has only {len(self.prs.slide_layouts)} layouts, expected 36+")
    
    def _load_config(self, config_path):
        """Load configuration or use defaults."""
        defaults = {
            'font_name': 'Outfit',
            'font_sizes': {
                'title': 24,
                'heading': 20,
                'body': 14,
                'caption': 12
            }
        }
        
        if config_path and Path(config_path).exists():
            with open(config_path) as f:
                user_config = yaml.safe_load(f)
                defaults.update(user_config)
        
        return defaults
    
    def _set_text_content(self, text_frame, content):
        """Set text content properly, handling bullets without duplication.
        
        Removes bullet characters (•) from YAML content since PowerPoint
        will add its own bullet formatting based on the template.
        """
        if not content:
            return
        
        # Clear existing text
        text_frame.clear()
        
        # Split content into lines
        lines = content.strip().split('\n')
        
        if not lines:
            return
        
        # Process lines
        current_paragraph = None
        
        for line in lines:
            original_line = line
            line = line.rstrip()  # Remove trailing whitespace but keep leading
            
            if not line.strip():
                # Empty line - add empty paragraph for spacing
                current_paragraph = text_frame.add_paragraph()
                current_paragraph.text = ""
                current_paragraph.level = 0
            elif line.strip().startswith('•') or line.strip().startswith('-'):
                # Bullet point - remove the bullet character (• or -) 
                # PowerPoint will add its own bullet formatting
                bullet_chars = ['•', '-', '*']
                cleaned_line = line.strip()
                for char in bullet_chars:
                    if cleaned_line.startswith(char):
                        cleaned_line = cleaned_line[1:].strip()
                        break
                
                current_paragraph = text_frame.add_paragraph()
                current_paragraph.text = cleaned_line
                current_paragraph.level = 0
            else:
                # Regular text line (section header or continuation)
                cleaned_line = line.strip()
                # Check if we should continue previous paragraph or start new one
                if current_paragraph and current_paragraph.text and not current_paragraph.text.endswith((':', '.', '!')):
                    # Continue previous paragraph with line break
                    current_paragraph.text += "\n" + cleaned_line
                else:
                    # Start new paragraph
                    current_paragraph = text_frame.add_paragraph()
                    current_paragraph.text = cleaned_line
                    current_paragraph.level = 0
    
    def generate(self, content_spec_path, output_path):
        """Generate presentation from content specification."""
        # Load content spec
        with open(content_spec_path) as f:
            spec = yaml.safe_load(f)
        
        # Clear template slides (keep only master)
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
        
        # Generate slides from spec
        for slide_spec in spec.get('slides', []):
            self._add_slide(slide_spec)
        
        # Save presentation
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_file))
        
        return {
            'success': True,
            'output': str(output_file),
            'slide_count': len(self.prs.slides)
        }
    
    def _add_slide(self, slide_spec):
        """Add a single slide based on specification."""
        # Get layout
        layout_type = slide_spec.get('type', 'title_only')
        layout_idx = LAYOUT_MAP.get(layout_type, 12)
        
        if layout_idx >= len(self.prs.slide_layouts):
            print(f"⚠️  Layout {layout_type} (index {layout_idx}) not found, using Title Only")
            layout_idx = 12
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        
        # Populate content based on layout type
        if layout_type in ['title_white', 'title_reverse']:
            self._populate_title_slide(slide, slide_spec)
        elif layout_type == 'divider':
            self._populate_divider(slide, slide_spec)
        elif layout_type == 'text_only':
            self._populate_text_only(slide, slide_spec)
        elif layout_type in ['two_column', 'three_column']:
            self._populate_columns(slide, slide_spec)
        elif layout_type == 'quote':
            self._populate_quote(slide, slide_spec)
        elif layout_type == 'end_slide':
            self._populate_end_slide(slide, slide_spec)
        else:
            # Default: populate title if present
            if hasattr(slide.shapes, 'title') and slide.shapes.title is not None and 'title' in slide_spec:
                slide.shapes.title.text = slide_spec['title']
        
        return slide
    
    def _populate_title_slide(self, slide, spec):
        """Populate title slide."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None and 'title' in spec:
            slide.shapes.title.text = spec['title']
        
        # Find subtitle placeholder (usually index 1)
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                shape.text = spec.get('subtitle', '')
                break
    
    def _populate_divider(self, slide, spec):
        """Populate section divider."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None and 'title' in spec:
            slide.shapes.title.text = spec['title']
    
    def _populate_text_only(self, slide, spec):
        """Populate text-only slide."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None and 'title' in spec:
            slide.shapes.title.text = spec['title']
        
        # Find content text box
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx > 0:
                    self._set_text_content(shape.text_frame, spec.get('content', ''))
                    break
    
    def _populate_columns(self, slide, spec):
        """Populate multi-column slide."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None and 'title' in spec:
            slide.shapes.title.text = spec['title']
        
        # Find column placeholders and populate
        columns = [spec.get('left_content', ''), 
                  spec.get('right_content', ''),
                  spec.get('middle_content', '')]
        
        col_idx = 0
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx > 0:
                    if col_idx < len(columns) and columns[col_idx]:
                        self._set_text_content(shape.text_frame, columns[col_idx])
                    col_idx += 1
    
    def _populate_quote(self, slide, spec):
        """Populate quote slide."""
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                if 'quote' in spec:
                    shape.text = spec['quote']
                if 'attribution' in spec:
                    # Add attribution to text frame
                    shape.text += f"\n\n— {spec['attribution']}"
                break
    
    def _populate_end_slide(self, slide, spec):
        """Populate end/thank you slide."""
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
            slide.shapes.title.text = spec.get('title', 'Thank You')
        
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx > 0:
                    self._set_text_content(shape.text_frame, spec.get('contact', ''))
                    break


def main():
    """CLI entry point."""
    if len(sys.argv) < 3:
        print("Usage: python3 ppt_generator.py <content_spec.yaml> <output.pptx>")
        print("\nExample:")
        print("  python3 ppt_generator.py input/content_spec.yaml output/presentation.pptx")
        sys.exit(1)
    
    content_spec = sys.argv[1]
    output_file = sys.argv[2]
    
    # Find template - try multiple locations
    script_dir = Path(__file__).parent.absolute()
    possible_template_paths = [
        script_dir.parent / "templates" / "HyFlux_Template_-.pptx",  # From scripts/ -> ../templates/
        script_dir / "templates" / "HyFlux_Template_-.pptx",  # scripts/templates/
        Path("templates") / "HyFlux_Template_-.pptx",  # Current directory
        Path("../templates") / "HyFlux_Template_-.pptx",  # One level up
    ]
    
    template = None
    for template_path in possible_template_paths:
        if template_path.exists():
            template = template_path
            break
    
    if not template:
        print(f"❌ Template not found: HyFlux_Template_-.pptx")
        print("   Searched in:")
        for path in possible_template_paths:
            print(f"     - {path}")
        print("   Place HyFlux_Template_-.pptx in templates/ directory")
        sys.exit(1)
    
    # Check content spec exists
    if not Path(content_spec).exists():
        print(f"❌ Content spec not found: {content_spec}")
        sys.exit(1)
    
    try:
        print(f"🚀 Generating presentation...")
        print(f"   Template: {template}")
        print(f"   Content:  {content_spec}")
        print(f"   Output:   {output_file}")
        
        generator = HyFluxPPTGenerator(str(template))
        result = generator.generate(content_spec, output_file)
        
        print(f"\n✅ Success!")
        print(f"   Created: {result['output']}")
        print(f"   Slides:  {result['slide_count']}")
        
    except Exception as e:
        print(f"\n❌ Generation failed: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
