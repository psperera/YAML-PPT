#!/usr/bin/env python3
"""
HyFlux PowerPoint Validator
Checks generated presentations for compliance with standards.
"""

import sys
from pathlib import Path
from pptx import Presentation
from collections import Counter

class HyFluxValidator:
    def __init__(self, pptx_path):
        self.path = Path(pptx_path)
        self.prs = Presentation(str(self.path))
        self.issues = []
        self.warnings = []
        self.info = []
    
    def validate_all(self):
        """Run all validation checks."""
        self._check_file_basics()
        self._check_dimensions()
        self._check_fonts()
        self._check_slide_count()
        self._check_placeholders()
        self._check_file_size()
        
        return self._generate_report()
    
    def _check_file_basics(self):
        """Validate file exists and is readable."""
        if not self.path.exists():
            self.issues.append(f"File not found: {self.path}")
            return False
        
        if not self.path.suffix == '.pptx':
            self.warnings.append(f"File extension is {self.path.suffix}, expected .pptx")
        
        self.info.append(f"✓ File readable: {self.path.name}")
        return True
    
    def _check_dimensions(self):
        """Check slide dimensions match 16:9 standard."""
        width = self.prs.slide_width
        height = self.prs.slide_height
        ratio = width / height
        
        expected_ratio = 16/9
        tolerance = 0.01
        
        if abs(ratio - expected_ratio) > tolerance:
            self.warnings.append(
                f"Aspect ratio is {ratio:.2f}:1, expected {expected_ratio:.2f}:1 (16:9)"
            )
        else:
            self.info.append(f"✓ Aspect ratio: 16:9")
        
        # Check dimensions
        expected_width = 12188825  # EMUs for 13.33 inches
        expected_height = 6858000   # EMUs for 7.5 inches
        
        if width == expected_width and height == expected_height:
            self.info.append(f"✓ Dimensions: 13.33\" × 7.5\"")
        else:
            actual_w_in = width / 914400
            actual_h_in = height / 914400
            self.warnings.append(
                f"Dimensions: {actual_w_in:.2f}\" × {actual_h_in:.2f}\", "
                f"expected 13.33\" × 7.5\""
            )
    
    def _check_fonts(self):
        """Check if Outfit font is used consistently."""
        fonts_used = set()
        font_counts = Counter()
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts_used.add(run.font.name)
                                font_counts[run.font.name] += 1
        
        outfit_fonts = {'Outfit', 'Outfit Semi Bold', 'Outfit SemiBold'}
        non_outfit = fonts_used - outfit_fonts
        
        if non_outfit:
            self.warnings.append(
                f"Non-Outfit fonts found: {', '.join(non_outfit)}"
            )
        else:
            self.info.append(f"✓ All fonts are Outfit family")
        
        # Show font distribution
        if fonts_used:
            self.info.append(f"  Fonts used: {', '.join(sorted(fonts_used))}")
    
    def _check_slide_count(self):
        """Check slide count is reasonable."""
        count = len(self.prs.slides)
        
        if count == 0:
            self.issues.append("Presentation has no slides")
        elif count > 100:
            self.warnings.append(
                f"Presentation has {count} slides (>100 may impact performance)"
            )
        else:
            self.info.append(f"✓ Slide count: {count}")
    
    def _check_placeholders(self):
        """Check for unfilled placeholders."""
        placeholder_texts = [
            'Click to add', 
            'Add text',
            '<insert',
            'placeholder',
            'TODO'
        ]
        
        found_placeholders = []
        
        for idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.lower()
                    for placeholder in placeholder_texts:
                        if placeholder.lower() in text:
                            found_placeholders.append(
                                f"Slide {idx}: '{shape.text[:50]}...'"
                            )
                            break
        
        if found_placeholders:
            self.warnings.append(
                f"Potential unfilled placeholders found:\n  " + 
                '\n  '.join(found_placeholders[:5])
            )
            if len(found_placeholders) > 5:
                self.warnings.append(f"  ... and {len(found_placeholders) - 5} more")
        else:
            self.info.append("✓ No obvious placeholders found")
    
    def _check_file_size(self):
        """Check file size is reasonable."""
        size_mb = self.path.stat().st_size / (1024 * 1024)
        
        if size_mb > 50:
            self.warnings.append(
                f"File size is {size_mb:.1f} MB (>50 MB may be slow to share)"
            )
        else:
            self.info.append(f"✓ File size: {size_mb:.1f} MB")
    
    def _generate_report(self):
        """Generate validation report."""
        print("\n" + "="*70)
        print("HyFlux PowerPoint Validation Report")
        print("="*70)
        
        if self.issues:
            print("\n❌ CRITICAL ISSUES:")
            for issue in self.issues:
                print(f"   • {issue}")
        
        if self.warnings:
            print("\n⚠️  WARNINGS:")
            for warning in self.warnings:
                print(f"   • {warning}")
        
        if self.info:
            print("\n✓ PASSED CHECKS:")
            for info in self.info:
                print(f"   {info}")
        
        print("\n" + "="*70)
        
        # Overall status
        if self.issues:
            print("Status: ❌ FAILED - Fix critical issues")
            return False
        elif self.warnings:
            print("Status: ⚠️  PASSED WITH WARNINGS - Review warnings")
            return True
        else:
            print("Status: ✅ PASSED - All checks OK")
            return True


def main():
    """CLI entry point."""
    if len(sys.argv) < 2:
        print("Usage: python3 validator.py <presentation.pptx>")
        print("\nExample:")
        print("  python3 validator.py output/presentation.pptx")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    try:
        validator = HyFluxValidator(pptx_file)
        passed = validator.validate_all()
        
        sys.exit(0 if passed else 1)
        
    except Exception as e:
        print(f"\n❌ Validation error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
